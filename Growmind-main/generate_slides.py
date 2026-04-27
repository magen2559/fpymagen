import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def apply_text_formatting(run, font_name="Calibri", font_size=Pt(18), color=None, bold=False):
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

def create_slide(prs, title_str, content_lines, is_title_slide=False, subtitle_str=""):
    if is_title_slide:
        slide_layout = prs.slide_layouts[0] # Title slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_str
        subtitle.text = subtitle_str
        
        # Format title
        for par in title.text_frame.paragraphs:
            par.alignment = PP_ALIGN.CENTER
            for run in par.runs:
                apply_text_formatting(run, font_name="Arial", font_size=Pt(44), color=RGBColor(0xA5, 0x2A, 0x2A), bold=True)
                
        # Format subtitle
        for par in subtitle.text_frame.paragraphs:
            par.alignment = PP_ALIGN.CENTER
            for run in par.runs:
                apply_text_formatting(run, font_name="Arial", font_size=Pt(20), color=RGBColor(0x33, 0x41, 0x55))
    else:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title.text = title_str
        for par in title.text_frame.paragraphs:
            for run in par.runs:
                apply_text_formatting(run, font_name="Arial", font_size=Pt(36), color=RGBColor(0xA5, 0x2A, 0x2A), bold=True)
        
        tf = body_shape.text_frame
        tf.text = content_lines[0] if content_lines else ""
        if content_lines:
            for run in tf.paragraphs[0].runs:
                apply_text_formatting(run, font_name="Arial", font_size=Pt(24), color=RGBColor(0x1E, 0x29, 0x3B))

        for line in content_lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
            # Check for sub-bullets
            if line.startswith("  -"):
                p.level = 1
                p.text = line.replace("  - ", "").strip()
            
            for run in p.runs:
                 apply_text_formatting(run, font_name="Arial", font_size=Pt(20) if p.level==1 else Pt(24), color=RGBColor(0x1E, 0x29, 0x3B))

def main():
    prs = Presentation()
    
    # Slide 1: Title
    create_slide(prs, "GrowMind", None, is_title_slide=True, 
                 subtitle_str="A Gamified Focus and Productivity App\nUsing a Pomodoro Timer and a Virtual 3D Garden\n\nBy: Izz Bin Ismail")
                 
    # Slide 2: Abstract
    create_slide(prs, "Abstract", [
        "Addresses student distraction and lack of focus.",
        "Integrates a Pomodoro timer with a reward system.",
        "Users build and nurture a virtual 3D garden through focused study.",
        "Positive reinforcement encourages sustained concentration.",
        "Gamifies the academic workflow and enhances student outcomes."
    ])
    
    # Slide 3: Problem Statement
    create_slide(prs, "Problem Statement", [
        "1. Delayed Gratification",
        "  - Academic success lacks immediate rewards.",
        "2. Digital Fragmentation",
        "  - Constant context-switching between study and entertainment.",
        "3. Academic Isolation",
        "  - Lack of community presence during solitary study sessions."
    ])

    # Slide 4: Objectives
    create_slide(prs, "Objectives", [
        "1. Develop a Pomodoro-style focus timer for study sessions.",
        "2. Implement an immersive virtual 3D garden that evolves.",
        "3. Create a coin-based reward system for successful focus cycles.",
        "4. Enable user authentication and progress tracking."
    ])

    # Slide 5: Methodology
    create_slide(prs, "Methodology", [
        "Agile Development Methodology Lifecycle:",
        "  - Planning: Defining scope and requirements",
        "  - Design: UI/UX, 3D Assets, System Architecture",
        "  - Implementation: Core modules, AI integration, 3D Engine",
        "  - Testing: Unit, Integration, Performance (60 FPS minimum)",
        "  - Deployment: Mobile app binaries (Expo)",
        "  - Analysis: Gathering feedback for future iterations"
    ])

    # Slide 6: Proposed Technical Stack
    create_slide(prs, "Technical Stack", [
        "Frontend: React Native + Expo",
        "3D/Graphics: React Three Fiber (R3F) via expo-gl",
        "Backend: Supabase (PostgreSQL, Auth)",
        "Logic & Services: Supabase Edge Functions (TypeScript/Deno)",
        "Realtime: Supabase Realtime for community visualization",
        "AI Engine: Google Gemini Pro via Edge Functions"
    ])

    # Slide 7: Core Features
    create_slide(prs, "Core Features", [
        "Focus Session Timer:",
        "  - Strict mode, tagging, and progress tracking.",
        "Interactive 3D Garden:",
        "  - Orbit controls, Procedural plant growth based on focus.",
        "AI Coach (GrowMind AI):",
        "  - Personalized productivity guidance powered by Gemini.",
        "Task Management & Shop:",
        "  - Customization with earned coins."
    ])

    # Slide 8: Commercialisation Potential
    create_slide(prs, "Commercialisation Potential", [
        "Freemium Models:",
        "  - Premium garden elements and exclusive themes.",
        "Institutional Licensing:",
        "  - Tailored solutions for universities and study centers.",
        "In-App Purchases:",
        "  - Specialized 3D assets and items.",
        "Data Insights:",
        "  - Analytics on student productivity patterns."
    ])

    # Slide 9: Conclusion
    create_slide(prs, "Conclusion & Future Work", [
        "Conclusion:",
        "  - Success in merging gamification with focus tools.",
        "  - Improved concentration and engagement.",
        "Future Work:",
        "  - Introduce deeper social features.",
        "  - Implement advanced analytics.",
        "  - Ensure cross-platform compatibility."
    ])

    # Slide 10: Q&A
    create_slide(prs, "Thank You", ["Questions?"])

    prs.save("c:/Users/daniel/Desktop/growmind/GrowMind_Presentation.pptx")
    print("Presentation created successfully at c:/Users/daniel/Desktop/growmind/GrowMind_Presentation.pptx")

if __name__ == "__main__":
    main()
